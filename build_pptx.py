"""Generate the project defense presentation as a .pptx file."""

from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parent
ASSETS = ROOT / "assets"
METRICS_PATH = ROOT / "reports" / "demo_run_v2" / "metrics.json"
OUT = ROOT / "reports" / "WebZipStudio_Presentation.pptx"

# Color palette: Midnight Executive
NAVY = RGBColor(0x1E, 0x27, 0x61)
ICE = RGBColor(0xCA, 0xDC, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x2E, 0x75, 0xB6)
DARK_TEXT = RGBColor(0x11, 0x18, 0x27)
MUTED = RGBColor(0x6B, 0x72, 0x80)


def load_metrics():
    if METRICS_PATH.exists():
        return json.loads(METRICS_PATH.read_text())
    return {}


METRICS = load_metrics()
TOTALS = METRICS.get("totals", {})
COMPARISON = METRICS.get("comparison", {})
BY_STRATEGY = METRICS.get("by_strategy", {})


# --- helpers ----------------------------------------------------
def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, x, y, w, h, text, *, font="Calibri",
                size=18, bold=False, color=DARK_TEXT,
                align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def add_bullets(slide, x, y, w, h, items, size=16, color=DARK_TEXT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = "•  " + item
        run.font.name = "Calibri"
        run.font.size = Pt(size)
        run.font.color.rgb = color
        p.space_after = Pt(6)
    return box


def add_stat_card(slide, x, y, w, h, value, label):
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.color.rgb = ICE
    bg.line.width = Pt(0.75)

    add_textbox(slide, x, y + Inches(0.15), w, Inches(0.6),
                value, size=26, bold=True, color=NAVY,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, x, y + Inches(0.75), w, Inches(0.4),
                label, size=11, color=MUTED, align=PP_ALIGN.CENTER)


# --- slides ----------------------------------------------------
def make_presentation() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Slide 1 — Title
    s = prs.slides.add_slide(blank)
    set_bg(s, NAVY)
    add_textbox(s, Inches(0.7), Inches(2.3), Inches(12.0), Inches(1.2),
                "WebZip Studio", size=58, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.7), Inches(3.4), Inches(12.0), Inches(0.7),
                "Webpage Compression System", size=28, color=ICE,
                align=PP_ALIGN.CENTER, italic=True)
    add_textbox(s, Inches(0.7), Inches(4.2), Inches(12.0), Inches(0.6),
                "Final Project — Data Structure and its Algorithms",
                size=18, color=ICE, align=PP_ALIGN.CENTER)
    members_text = "Bakbergen Amir 202469990559   ·   Bakbergen Alen 202469990562   ·   Huang Liu Diego David 202469990549"
    add_textbox(s, Inches(0.7), Inches(6.4), Inches(12.0), Inches(0.5),
                members_text, size=14, color=ICE, align=PP_ALIGN.CENTER)

    # Slide 2 — Why this project
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    add_textbox(s, Inches(0.7), Inches(0.5), Inches(12.0), Inches(0.8),
                "Why a webpage compression system", size=34, bold=True, color=NAVY)
    add_bullets(s, Inches(0.9), Inches(1.6), Inches(11.5), Inches(5.0), [
        "Webpages keep growing — landing pages routinely ship megabytes of HTML, CSS, JS, and images.",
        "Generic tools like gzip and zip work, but they hide the data structures that the course is about.",
        "We built a layered system that exposes every algorithm: LZ77 + Huffman for text, Pillow for images.",
        "Output is packaged with a manifest, integrity hashes, and metrics — defensible end-to-end.",
        "Every claim in the report can be re-derived by running tests/run_all.py and the GUI itself.",
    ], size=18)

    # Slide 3 — Goals
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    add_textbox(s, Inches(0.7), Inches(0.5), Inches(12.0), Inches(0.8),
                "What the system delivers", size=34, bold=True, color=NAVY)
    cards = [
        ("HTML / CSS / JS", "Lossless LZ77 + Huffman"),
        ("JPG / PNG", "Quality presets (High / Balanced / Strong)"),
        ("Auto strategy", "Picked from file extension"),
        ("Decompress + verify", "SHA-256 MATCH / MISMATCH"),
        ("Compare", "vs gzip and ZIP_DEFLATED"),
        ("Transfer simulation", "4G / 5G / WiFi"),
        ("Visualizer", "Top tokens, codes, LZ77 matches"),
        ("Incremental", "Skip files whose hash hasn't changed"),
    ]
    cw, ch = Inches(2.85), Inches(1.4)
    for idx, (val, lab) in enumerate(cards):
        col = idx % 4
        row = idx // 4
        x = Inches(0.7 + col * 3.1)
        y = Inches(1.7 + row * 1.7)
        add_stat_card(s, x, y, cw, ch, val, lab)

    # Slide 4 — Architecture
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    add_textbox(s, Inches(0.7), Inches(0.5), Inches(12.0), Inches(0.8),
                "Architecture", size=34, bold=True, color=NAVY)
    pipeline = [
        "1. User selects files / folder",
        "2. StrategySelector (hash map: extension -> strategy)",
        "3. TextPipeline   ➜   LZ77 sliding window  ➜  Huffman tree (heap + binary tree)  ➜  .wzs",
        "4. ImagePipeline  ➜   Pillow re-encode at quality preset",
        "5. ManifestManager + IntegrityVerifier + MetricsCollector",
        "6. Decompress reverses every step; SHA-256 confirms exact restoration",
    ]
    add_bullets(s, Inches(0.9), Inches(1.6), Inches(11.5), Inches(5.0),
                pipeline, size=18)

    # Slide 5 — Data structures
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    add_textbox(s, Inches(0.7), Inches(0.5), Inches(12.0), Inches(0.8),
                "Data structures in the system", size=34, bold=True, color=NAVY)
    ds_rows = [
        ("Hash map", "Frequency table; LZ77 prefix index; extension router; incremental cache"),
        ("Min-heap (priority queue)", "Build the Huffman tree from least-frequent symbols upward"),
        ("Binary tree", "Huffman code derivation via DFS"),
        ("Sliding window", "LZ77 lookback buffer; bounded chain length keeps it linear in practice"),
        ("Queue (deque)", "FIFO file processing in CompressionManager"),
        ("Set", "De-duplicate dropped paths; diff old vs new file hashes"),
        ("List", "Token streams; per-file metric tables; chart series"),
    ]
    table = s.shapes.add_table(rows=len(ds_rows) + 1, cols=2,
                               left=Inches(0.7), top=Inches(1.6),
                               width=Inches(11.9), height=Inches(5.0)).table
    table.columns[0].width = Inches(3.5)
    table.columns[1].width = Inches(8.4)
    h_cells = table.rows[0].cells
    for i, h in enumerate(["Structure", "Where it lives"]):
        h_cells[i].text = h
        for p in h_cells[i].text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True
                r.font.color.rgb = WHITE
                r.font.size = Pt(14)
        h_cells[i].fill.solid()
        h_cells[i].fill.fore_color.rgb = NAVY
    for i, (k, v) in enumerate(ds_rows, start=1):
        cells = table.rows[i].cells
        cells[0].text = k
        cells[1].text = v
        for c in cells:
            for p in c.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(13)

    # Slide 6 — Text pipeline
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    add_textbox(s, Inches(0.7), Inches(0.5), Inches(12.0), Inches(0.8),
                "Text pipeline: LZ77 + Huffman", size=34, bold=True, color=NAVY)
    add_bullets(s, Inches(0.9), Inches(1.6), Inches(11.5), Inches(2.4), [
        "LZ77 emits (offset, length, next_byte) triples by walking a sliding window.",
        "A 3-byte hash-map prefix index speeds up match search; chains are capped at 64.",
        "Huffman builds a min-heap, then a binary tree, then a prefix code book.",
        "The bitstream is packed into a .wzs container with a JSON header listing frequencies.",
    ], size=18)
    add_textbox(s, Inches(0.9), Inches(4.2), Inches(11.5), Inches(0.5),
                ".wzs format", size=18, bold=True, color=NAVY)
    add_textbox(s, Inches(0.9), Inches(4.7), Inches(11.5), Inches(0.5),
                "magic 'WZS1' (4B)  |  uint32 header_len  |  JSON header  |  packed Huffman bits",
                size=14, color=DARK_TEXT)
    add_textbox(s, Inches(0.9), Inches(5.6), Inches(11.5), Inches(1.2),
                "Decompression reverses every step. Text round-trip is verified by SHA-256.",
                size=14, color=MUTED, italic=True)

    # Slide 7 — Image pipeline
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    add_textbox(s, Inches(0.7), Inches(0.5), Inches(12.0), Inches(0.8),
                "Image pipeline (Pillow + quality presets)", size=34, bold=True, color=NAVY)
    add_bullets(s, Inches(0.9), Inches(1.6), Inches(11.5), Inches(5.0), [
        "JPG: re-encode with quality 90 (High) / 75 (Balanced) / 55 (Strong); progressive + optimize.",
        "PNG: optimize + optional palette reduction in Strong mode (skipped if image has alpha).",
        "Compression is lossy by design — we expose the trade-off rather than hide it.",
        "Originals are never silently re-encoded into a different format.",
        "Integrity check applies only to text; the image case shows OK / MISSING.",
    ], size=18)

    # Slide 8 — Demo numbers (stats)
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    add_textbox(s, Inches(0.7), Inches(0.5), Inches(12.0), Inches(0.8),
                "Demo run — sample_large_case (16 files)", size=32, bold=True, color=NAVY)
    cards = [
        (str(int(TOTALS.get("file_count", 16))), "Files"),
        (f"{TOTALS.get('original_size', 0)/1024:.1f} KB", "Original"),
        (f"{TOTALS.get('compressed_size', 0)/1024:.1f} KB", "Compressed"),
        (f"{TOTALS.get('ratio', 1.0):.3f}", "Ratio"),
        (f"{TOTALS.get('savings_pct', 0.0):.1f}%", "Savings"),
    ]
    cw, ch = Inches(2.3), Inches(1.4)
    for idx, (val, lab) in enumerate(cards):
        x = Inches(0.7 + idx * 2.4)
        add_stat_card(s, x, Inches(1.7), cw, ch, val, lab)

    # Bullet detail
    add_bullets(s, Inches(0.9), Inches(3.5), Inches(11.5), Inches(3.5), [
        f"Custom .wzs total: {COMPARISON.get('custom', {}).get('compressed_size', 0)/1024:.2f} KB "
        f"in {COMPARISON.get('custom', {}).get('duration_seconds', 0.0)*1000:.1f} ms",
        f"gzip total: {COMPARISON.get('gzip', {}).get('compressed_size', 0)/1024:.2f} KB "
        f"in {COMPARISON.get('gzip', {}).get('duration_seconds', 0.0)*1000:.1f} ms",
        f"zip total:  {COMPARISON.get('zip', {}).get('compressed_size', 0)/1024:.2f} KB "
        f"in {COMPARISON.get('zip', {}).get('duration_seconds', 0.0)*1000:.1f} ms",
        "gzip wins on raw ratio because it ships pre-baked Huffman tables — that is the trade-off the report calls out.",
    ], size=16)

    # Slide 9 — Comparison chart
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    add_textbox(s, Inches(0.7), Inches(0.5), Inches(12.0), Inches(0.8),
                "Custom vs gzip vs zip", size=34, bold=True, color=NAVY)
    chart_path = ASSETS / "chart_compare.png"
    if chart_path.exists():
        s.shapes.add_picture(str(chart_path), Inches(2.0), Inches(1.4),
                             width=Inches(9.3))

    # Slide 10 — Per-strategy chart
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    add_textbox(s, Inches(0.7), Inches(0.5), Inches(12.0), Inches(0.8),
                "Bytes saved by strategy", size=34, bold=True, color=NAVY)
    chart_path = ASSETS / "chart_strategy.png"
    if chart_path.exists():
        s.shapes.add_picture(str(chart_path), Inches(3.5), Inches(1.4),
                             height=Inches(5.4))

    # Slide 11 — Per-file chart
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    add_textbox(s, Inches(0.7), Inches(0.5), Inches(12.0), Inches(0.8),
                "Per-file savings", size=34, bold=True, color=NAVY)
    chart_path = ASSETS / "chart_per_file.png"
    if chart_path.exists():
        s.shapes.add_picture(str(chart_path), Inches(2.0), Inches(1.4),
                             width=Inches(9.3))

    # Slide 12 — Complexity
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    add_textbox(s, Inches(0.7), Inches(0.5), Inches(12.0), Inches(0.8),
                "Complexity at a glance", size=34, bold=True, color=NAVY)
    rows = [
        ("Frequency counting", "O(n) time, O(k) space"),
        ("Huffman tree (heap)", "O(k log k) time, O(k) space"),
        ("Huffman encode/decode", "O(n) time, O(n) space"),
        ("LZ77 with hash index", "O(n · c · ℓ) ≈ O(n) for fixed c, ℓ"),
        ("LZ77 decode", "O(n)"),
        ("Whole batch", "Linear in total bytes"),
    ]
    table = s.shapes.add_table(rows=len(rows) + 1, cols=2,
                               left=Inches(0.9), top=Inches(1.6),
                               width=Inches(11.5), height=Inches(5.0)).table
    table.columns[0].width = Inches(4.0)
    table.columns[1].width = Inches(7.5)
    h_cells = table.rows[0].cells
    for i, h in enumerate(["Stage", "Big-O"]):
        h_cells[i].text = h
        for p in h_cells[i].text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True
                r.font.color.rgb = WHITE
                r.font.size = Pt(14)
        h_cells[i].fill.solid()
        h_cells[i].fill.fore_color.rgb = NAVY
    for i, (k, v) in enumerate(rows, start=1):
        cells = table.rows[i].cells
        cells[0].text = k
        cells[1].text = v
        for c in cells:
            for p in c.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(14)

    # Slide 13 — GUI tour
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    add_textbox(s, Inches(0.7), Inches(0.5), Inches(12.0), Inches(0.8),
                "GUI tour (5 tabs)", size=34, bold=True, color=NAVY)
    add_bullets(s, Inches(0.9), Inches(1.6), Inches(11.5), Inches(5.0), [
        "Compress — drag-and-drop, file list, output folder, quality preset, incremental toggle.",
        "Decompress — pick a package; restored files turn green when SHA-256 matches the manifest.",
        "Analytics — stat cards, gzip/zip comparison, transfer-time table, bar chart.",
        "Visualizer — top tokens, Huffman code book, sample LZ77 matches.",
        "Incremental — UNCHANGED / RECOMPRESSED / ADDED / REMOVED buckets per file.",
    ], size=18)

    # Slide 14 — Contributions
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    add_textbox(s, Inches(0.7), Inches(0.5), Inches(12.0), Inches(0.8),
                "Who did what", size=34, bold=True, color=NAVY)
    rows = [
        ("Bakbergen Amir\n202469990559",
         "Lead — Huffman + LZ77 + text pipeline + manifest + tests"),
        ("Bakbergen Alen\n202469990562",
         "GUI (PySide6, 5 tabs) + workers + comparison + transfer + chart"),
        ("Huang Liu Diego David\n202469990549",
         "Image pipeline + sample data + benchmarks + report charts"),
    ]
    table = s.shapes.add_table(rows=len(rows) + 1, cols=2,
                               left=Inches(0.9), top=Inches(1.6),
                               width=Inches(11.5), height=Inches(4.0)).table
    table.columns[0].width = Inches(4.0)
    table.columns[1].width = Inches(7.5)
    h = table.rows[0].cells
    for i, t in enumerate(["Member", "Owns"]):
        h[i].text = t
        for p in h[i].text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True
                r.font.color.rgb = WHITE
                r.font.size = Pt(14)
        h[i].fill.solid()
        h[i].fill.fore_color.rgb = NAVY
    for i, (k, v) in enumerate(rows, start=1):
        cells = table.rows[i].cells
        cells[0].text = k
        cells[1].text = v
        for c in cells:
            for p in c.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(13)

    # Slide 15 — Q&A
    s = prs.slides.add_slide(blank)
    set_bg(s, NAVY)
    add_textbox(s, Inches(0.7), Inches(2.4), Inches(12.0), Inches(1.4),
                "Thank you", size=64, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.7), Inches(3.8), Inches(12.0), Inches(0.8),
                "Questions?", size=28, color=ICE,
                align=PP_ALIGN.CENTER, italic=True)
    add_textbox(s, Inches(0.7), Inches(6.5), Inches(12.0), Inches(0.5),
                "WebZip Studio — Data Structure and its Algorithms",
                size=14, color=ICE, align=PP_ALIGN.CENTER)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    out = make_presentation()
    print(f"Wrote {out}")
